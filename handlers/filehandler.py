import os
import asyncio
import tempfile
import logging
import re
from typing import List, Optional

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

class FileHandler:
    def __init__(self):
        os.makedirs("temp", exist_ok=True)

    async def process_attachments(self, attachments: List) -> str:
        file_contexts = []
        for attachment in attachments:
            filename = attachment.filename.lower()
            if filename.endswith('.pdf'):
                content = await self.extract_pdf(attachment)
                if content:
                    file_contexts.append(f"[PDF: {attachment.filename}]\n{content}")
            elif filename.endswith('.docx'):
                content = await self.extract_docx(attachment)
                if content:
                    file_contexts.append(f"[DOCX: {attachment.filename}]\n{content}")
            elif filename.endswith('.pptx'):
                content = await self.extract_pptx(attachment)
                if content:
                    file_contexts.append(f"[PPTX: {attachment.filename}]\n{content}")
            elif filename.endswith('.txt'):
                content = await self.extract_txt(attachment)
                if content:
                    file_contexts.append(f"[TXT: {attachment.filename}]\n{content}")
        return "\n".join(file_contexts)

    async def extract_pdf(self, attachment) -> Optional[str]:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                await attachment.save(tmp.name)
                file_path = tmp.name

            def read_pdf(path):
                with fitz.open(path) as doc:
                    return "".join(page.get_text() for page in doc)

            text = await asyncio.to_thread(read_pdf, file_path)
            os.remove(file_path)
            return self.clean_text(text) if text else None
        except Exception as e:
            logger.exception(f"Error reading PDF: {e}")
            return f"[Error reading PDF: {e}]"

    async def extract_docx(self, attachment) -> Optional[str]:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                await attachment.save(tmp.name)
                file_path = tmp.name

            def read_docx(path):
                doc = Document(path)
                return "\n".join(para.text for para in doc.paragraphs)

            text = await asyncio.to_thread(read_docx, file_path)
            os.remove(file_path)
            return self.clean_text(text) if text else None
        except Exception as e:
            logger.exception(f"Error reading DOCX: {e}")
            return f"[Error reading DOCX: {e}]"

    async def extract_txt(self, attachment) -> Optional[str]:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".txt", mode="w+b") as tmp:
                await attachment.save(tmp.name)
                file_path = tmp.name

            def read_txt(path):
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()

            text = await asyncio.to_thread(read_txt, file_path)
            os.remove(file_path)
            return self.clean_text(text) if text else None
        except Exception as e:
            logger.exception(f"Error reading TXT: {e}")
            return f"[Error reading TXT: {e}]"

    async def extract_pptx(self, attachment) -> Optional[str]:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                await attachment.save(tmp.name)
                file_path = tmp.name

            def read_pptx(path):
                prs = Presentation(path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                            for sub_shape in shape.shapes:
                                if hasattr(sub_shape, "text"):
                                    text += sub_shape.text + "\n"
                        elif hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text

            text = await asyncio.to_thread(read_pptx, file_path)
            os.remove(file_path)
            return self.clean_text(text) if text else None
        except Exception as e:
            logger.exception(f"Error reading PPTX: {e}")
            return f"[Error reading PPTX: {e}]"

    def clean_text(self, text: str) -> str:
        """Normalize whitespace and repeated newlines."""
        return re.sub(r'\n\s*\n+', '\n\n', text.strip())
